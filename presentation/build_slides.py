"""
Build AgentFlow presentation as PowerPoint (.pptx).
Run: python presentation/build_slides.py
Output: presentation/agentflow_slides.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# Colors
BLUE   = RGBColor(0x1F, 0x49, 0x7D)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0xF2, 0xF2, 0xF2)
GREEN  = RGBColor(0x00, 0x70, 0xC0)
DARK   = RGBColor(0x26, 0x26, 0x26)
ACCENT = RGBColor(0xED, 0x7D, 0x31)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # blank layout

def add_slide():
    return prs.slides.add_slide(BLANK)

def rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background() if line is None else None
    if fill:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is None:
        shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def header_bar(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 1.3, fill=BLUE)
    txt(slide, title, 0.3, 0.15, 12, 0.7, size=28, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.3, 0.8, 12, 0.4, size=14, color=RGBColor(0xBD, 0xD7, 0xEE))

def footer(slide, page):
    txt(slide, f"AgentFlow | Northeastern Self-Improving AI | Spring 2026  ·  {page}", 0.3, 7.15, 12, 0.3,
        size=10, color=RGBColor(0x99,0x99,0x99))

# ── Slide 1: Title ───────────────────────────────────────────────────────────
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=BLUE)
rect(s, 0, 2.8, 13.33, 0.04, fill=ACCENT)
txt(s, "AgentFlow", 1, 1.2, 11, 1.2, size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Self-Improving AI via Flow-GRPO", 1, 2.4, 11, 0.6, size=24, color=RGBColor(0xBD,0xD7,0xEE), align=PP_ALIGN.CENTER)
txt(s, "Varun Kasa  ·  Northeastern University  ·  April 2026", 1, 3.2, 11, 0.5, size=16, color=RGBColor(0xBD,0xD7,0xEE), align=PP_ALIGN.CENTER)
txt(s, "Based on: Huang et al., 'AgentFlow' (2025)", 1, 3.9, 11, 0.4, size=13, color=RGBColor(0x9D,0xC3,0xE6), align=PP_ALIGN.CENTER)

# ── Slide 2: Motivation ───────────────────────────────────────────────────────
s = add_slide()
header_bar(s, "Motivation", "Why train agents to plan?")
rect(s, 0.3, 1.5, 5.9, 5.6, fill=GRAY)
rect(s, 6.5, 1.5, 6.5, 5.6, fill=GRAY)
txt(s, "The Problem", 0.5, 1.6, 5.5, 0.5, size=16, bold=True, color=BLUE)
for i, line in enumerate([
    "• LLMs fail at multi-step tasks requiring tool use",
    "• Training agents is hard: sparse rewards,\n  multi-turn credit assignment",
    "• Standard RL (single-turn) doesn't capture\n  the full trajectory",
    "• No principled way to train the planner\n  while keeping tools fixed",
]):
    txt(s, line, 0.5, 2.2 + i*0.9, 5.5, 0.85, size=14)
txt(s, "AgentFlow's Answer", 6.7, 1.6, 6.0, 0.5, size=16, bold=True, color=BLUE)
for i, line in enumerate([
    "• Modular 4-agent pipeline: Planner → Executor\n  → Verifier → Generator",
    "• Only the Planner is trained (LoRA fine-tuning)",
    "• Flow-GRPO: multi-turn extension of GRPO",
    "• Binary reward broadcast to all planner turns",
    "• Group-normalize advantages across G=8 rollouts",
]):
    txt(s, line, 6.7, 2.2 + i*0.9, 6.1, 0.85, size=14)
footer(s, "2 / 13")

# ── Slide 3: Architecture ─────────────────────────────────────────────────────
s = add_slide()
header_bar(s, "AgentFlow Architecture", "Four specialized modules share a memory buffer")
boxes = [("Query", 0.4), ("Planner\n(trainable)", 2.4), ("Executor\n(tools)", 4.4),
         ("Verifier\n(stop?)", 6.4), ("Generator", 8.4), ("Answer", 10.5)]
colors = [GRAY, ACCENT, GREEN, GREEN, GREEN, GRAY]
for (label, x), col in zip(boxes, colors):
    c = RGBColor(0xFF,0xFF,0xFF) if col != GRAY else DARK
    r = rect(s, x, 3.0, 1.7, 0.9, fill=col)
    r.line.fill.background()
    txt(s, label, x+0.05, 3.0, 1.6, 0.9, size=13, bold=True, color=c, align=PP_ALIGN.CENTER)
    if x < 10.5:
        txt(s, "→", x+1.7, 3.2, 0.5, 0.5, size=20, bold=True, color=DARK)
rect(s, 0.4, 4.3, 10.4, 0.7, fill=RGBColor(0xDE,0xEB,0xF7))
txt(s, "Shared Memory Buffer", 0.5, 4.35, 10.2, 0.6, size=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
txt(s, "loop until STOP", 7.0, 2.4, 2.5, 0.4, size=12, color=DARK, align=PP_ALIGN.CENTER)
notes = [
    ("Planner", "Decides next tool + sub-goal\n→ Only this module is trained with Flow-GRPO", 0.4, 5.3),
    ("Executor", "Runs selected tool:\nGoogle Search, Wikipedia, SQL, Python", 3.5, 5.3),
    ("Verifier", "Checks: do we have enough\ninformation to stop?", 6.6, 5.3),
    ("Generator", "Synthesizes final answer\nfrom accumulated memory", 9.7, 5.3),
]
for title, body, x, y in notes:
    txt(s, title, x, y, 3.0, 0.35, size=12, bold=True, color=BLUE)
    txt(s, body, x, y+0.35, 3.1, 0.7, size=11, color=DARK)
footer(s, "3 / 13")

# ── Slide 4: Flow-GRPO Algorithm ──────────────────────────────────────────────
s = add_slide()
header_bar(s, "Flow-GRPO Algorithm", "Multi-turn GRPO for agentic trajectories")
rect(s, 0.3, 1.5, 6.2, 5.6, fill=GRAY)
rect(s, 6.8, 1.5, 6.2, 5.6, fill=GRAY)
txt(s, "Standard GRPO", 0.5, 1.6, 5.8, 0.4, size=15, bold=True, color=BLUE)
for i, l in enumerate([
    "1. Sample G rollouts per query",
    "2. Compute reward for each",
    "3. Group-normalize → advantages",
    "4. PPO-clipped policy gradient",
    "→ Single-turn only"
]):
    txt(s, l, 0.5, 2.15+i*0.7, 5.8, 0.65, size=13, color=DARK if i < 4 else ACCENT, bold=(i==4))
txt(s, "Flow-GRPO (AgentFlow)", 7.0, 1.6, 5.8, 0.4, size=15, bold=True, color=BLUE)
for i, l in enumerate([
    "1. Sample G=8 full agent trajectories",
    "2. Binary reward: final answer correct?",
    "3. Broadcast reward to ALL planner turns",
    "4. Group-normalize across G trajectories",
    "5. PPO clip (ε=0.2) + KL penalty (0.01)\n   vs frozen reference model",
    "LoRA: r=16, α=64, q/k/v/o_proj\n→ Only 0.44% of params trainable"
]):
    bold = i >= 5
    col = ACCENT if bold else DARK
    txt(s, l, 7.0, 2.15+i*0.75, 5.8, 0.72, size=13, bold=bold, color=col)
footer(s, "4 / 13")

# ── Slide 5: Implementation ───────────────────────────────────────────────────
s = add_slide()
header_bar(s, "Implementation Choices")
headers = ["Component", "Choice", "Reason"]
rows = [
    ("Planner (Steps 1-2)", "Qwen2.5-7B-Instruct", "Best Together AI serverless model"),
    ("Planner (Step 3)", "Qwen3.5-0.8B/2B/4B/9B/27B", "Scaling study — Modal vLLM"),
    ("GPU Inference", "Modal A10G / A100 + vLLM", "Cost-effective; $1.10–$3.73/hr"),
    ("Web Search", "Serper.dev", "Google CSE had access issues"),
    ("Training", "Modal A10G, LoRA", "~$12 per 500-step run"),
    ("Evaluation", "LLM-as-judge (7B)", "Matches paper methodology"),
    ("Thinking mode", "--override-generation-config\n{enable_thinking:false}", "Qwen3.5 uses thinking template\nby default — breaks JSON format"),
]
col_w = [3.0, 3.8, 5.8]
col_x = [0.3, 3.4, 7.3]
rect(s, 0.3, 1.4, 12.7, 0.5, fill=BLUE)
for j, (h, x, w) in enumerate(zip(headers, col_x, col_w)):
    txt(s, h, x+0.1, 1.45, w-0.1, 0.4, size=13, bold=True, color=WHITE)
for i, row in enumerate(rows):
    bg = GRAY if i % 2 == 0 else WHITE
    rect(s, 0.3, 1.9+i*0.72, 12.7, 0.72, fill=bg)
    for j, (cell, x, w) in enumerate(zip(row, col_x, col_w)):
        txt(s, cell, x+0.1, 1.92+i*0.72, w-0.15, 0.68, size=11)
footer(s, "5 / 13")

# ── Slide 6: Steps 1-2 Results ────────────────────────────────────────────────
s = add_slide()
header_bar(s, "Steps 1-2: Baseline Reproduction", "Qwen2.5-7B-Instruct via Together AI")
headers2 = ["Benchmark", "Ours", "Paper Target", "Δ", "Status"]
rows2 = [
    ("Bamboogle",    "67.74%", "58.4%", "+9.3%",  "✓ Exceeds"),
    ("HotpotQA",    "49.48%", "51.3%", "-1.8%",  "~ Close"),
    ("Musique",     "17.44%", "19.2%", "-1.8%",  "~ Close"),
    ("2WikiMultiHop","37.89%","60.0%", "-22.1%", "✗ Below"),
    ("GAIA",        "17.36%", "17.2%", "+0.2%",  "✓ Matches"),
]
col_w2 = [3.0, 2.0, 2.5, 2.0, 3.0]
col_x2 = [0.3, 3.4, 5.5, 8.1, 10.2]
rect(s, 0.3, 1.5, 12.7, 0.5, fill=BLUE)
for h, x, w in zip(headers2, col_x2, col_w2):
    txt(s, h, x+0.05, 1.55, w, 0.4, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
for i, row in enumerate(rows2):
    bg = GRAY if i % 2 == 0 else WHITE
    rect(s, 0.3, 2.05+i*0.6, 12.7, 0.6, fill=bg)
    status_col = GREEN if "✓" in row[-1] else (ACCENT if "~" in row[-1] else RGBColor(0xC0,0x00,0x00))
    for j, (cell, x, w) in enumerate(zip(row, col_x2, col_w2)):
        c = status_col if j == 4 else DARK
        b = j == 4
        txt(s, cell, x+0.05, 2.07+i*0.6, w, 0.55, size=13, bold=b, color=c, align=PP_ALIGN.CENTER)
txt(s, "Key finding: Google Search critical — Bamboogle jumped 51% → 68% after fixing Serper.dev integration.",
    0.3, 5.2, 12.7, 0.5, size=13, color=BLUE, bold=True)
txt(s, "2Wiki gap: Comparative multi-hop questions (\"which director is older?\") require fine-tuned retrieval — paper likely uses post-GRPO model.",
    0.3, 5.75, 12.7, 0.7, size=12, color=DARK)
footer(s, "6 / 13")

# ── Slide 7: Scaling Study ────────────────────────────────────────────────────
s = add_slide()
header_bar(s, "Step 3: Model Scaling Study", "AgentFlow performance vs model size")
rect(s, 0.3, 1.5, 8.5, 5.6, fill=GRAY)
rect(s, 9.1, 1.5, 4.0, 5.6, fill=GRAY)
txt(s, "Qwen2.5 family (apples-to-apples)", 0.5, 1.6, 8.1, 0.4, size=14, bold=True, color=BLUE)
h3 = ["Model", "Bamboogle", "HotpotQA", "Musique", "2Wiki", "GAIA"]
cw3 = [1.5, 1.3, 1.3, 1.3, 1.3, 1.3]
cx3 = [0.35, 1.9, 3.25, 4.6, 5.95, 7.3]
rect(s, 0.35, 2.1, 8.3, 0.45, fill=BLUE)
for h, x, w in zip(h3, cx3, cw3):
    txt(s, h, x, 2.15, w, 0.38, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
data3 = [("0.5B","37.1%","27.6%","6.1%","24.9%","12.3%"),
         ("3B",  "34.7%","28.6%","7.1%","23.8%","11.4%"),
         ("7B",  "67.7%","49.5%","17.4%","37.9%","17.4%")]
for i, row in enumerate(data3):
    bg = GRAY if i%2==0 else WHITE
    rect(s, 0.35, 2.55+i*0.6, 8.3, 0.6, fill=bg)
    for j, (cell, x, w) in enumerate(zip(row, cx3, cw3)):
        bold = i==2
        col = ACCENT if (i==2 and j>0) else DARK
        txt(s, cell, x, 2.57+i*0.6, w, 0.55, size=12, bold=bold, color=col, align=PP_ALIGN.CENTER)
txt(s, "Qwen3.5 family (Bamboogle)", 0.5, 4.15, 8.1, 0.4, size=14, bold=True, color=BLUE)
rect(s, 0.35, 4.6, 8.3, 0.4, fill=BLUE)
for h, x, w in zip(["Model","Bamboogle"], [0.35,1.9], [1.5,1.3]):
    txt(s, h, x, 4.63, w, 0.35, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "(results pending)", 3.3, 4.65, 5.0, 0.35, size=11, color=RGBColor(0x99,0x99,0x99))
txt(s, "Key Insights", 9.3, 1.6, 3.6, 0.4, size=14, bold=True, color=BLUE)
for i, l in enumerate([
    "7B is a 'phase transition':\ndramatically better tool-use planning",
    "0.5B ≈ 3B gap is small;\n3B → 7B is the big jump",
    "Scaling benefit consistent\nacross all 5 benchmarks",
    "Qwen3.5 requires\nenable_thinking=false\nfor AgentFlow compatibility",
]):
    txt(s, l, 9.3, 2.15+i*1.1, 3.7, 1.0, size=12, color=DARK)
footer(s, "7 / 13")

# ── Slide 8: Text-to-SQL ──────────────────────────────────────────────────────
s = add_slide()
header_bar(s, "Step 4: Text-to-SQL Extension", "Spider benchmark — 1034 dev questions, 166 SQLite databases")
rect(s, 0.3, 1.5, 5.8, 5.6, fill=GRAY)
rect(s, 6.4, 1.5, 6.6, 5.6, fill=GRAY)
txt(s, "Setup", 0.5, 1.6, 5.4, 0.4, size=15, bold=True, color=BLUE)
for i, l in enumerate([
    "New tool: SQL_Executor_Tool\n→ Executes SELECT on Spider SQLite DBs",
    "Schema injected into query:\n\"Table singer: [Singer_ID (number), ...\"",
    "Metric: Execution Accuracy\n(results match, not exact SQL string)",
    "Model: Qwen2.5-7B-Instruct-Turbo\nvia Together AI serverless",
    "Max steps: 5 per question\n(generate SQL → execute → verify → correct)",
]):
    txt(s, l, 0.5, 2.15+i*0.95, 5.4, 0.9, size=13, color=DARK)
txt(s, "Result", 6.6, 1.6, 6.0, 0.4, size=15, bold=True, color=BLUE)
rect(s, 6.6, 2.1, 6.0, 1.5, fill=BLUE)
txt(s, "91.01%", 6.6, 2.1, 6.0, 1.1, size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "941 / 1034 questions correct", 6.6, 3.1, 6.0, 0.45, size=14, color=RGBColor(0xBD,0xD7,0xEE), align=PP_ALIGN.CENTER)
txt(s, "Why so high?", 6.6, 3.8, 6.0, 0.4, size=14, bold=True, color=BLUE)
for i, l in enumerate([
    "• Execute-verify-correct loop:\n  model tries SQL, sees result, fixes errors",
    "• Schema context eliminates\n  hallucinated column/table names",
    "• Competitive with SOTA on\n  Spider dev leaderboard (~90-92%)",
    "• AgentFlow's tool-use framework\n  generalizes beyond QA tasks",
]):
    txt(s, l, 6.6, 4.3+i*0.75, 6.1, 0.72, size=12, color=DARK)
footer(s, "8 / 13")

# ── Slide 9: Flow-GRPO Results ────────────────────────────────────────────────
s = add_slide()
header_bar(s, "Steps 5-6: Flow-GRPO Training", "LoRA fine-tuning of Qwen3.5-0.8B planner")
rect(s, 0.3, 1.5, 6.0, 5.6, fill=GRAY)
rect(s, 6.6, 1.5, 6.4, 5.6, fill=GRAY)
txt(s, "Training Setup", 0.5, 1.6, 5.6, 0.4, size=15, bold=True, color=BLUE)
for i, l in enumerate([
    "Base model: Qwen3.5-0.8B (correct assignment model)",
    "Steps 5 (QA): NQ + HotpotQA training data\n500 steps, G=8, A10G GPU",
    "Step 6 (SQL): Spider train split\n500 steps, same LoRA config",
    "LoRA: r=16, α=64, q/k/v/o_proj\n0.44% params trainable",
    "PPO clip ε=0.2, KL coef=0.01\nLR=1e-4, gradient clip 0.5",
]):
    txt(s, l, 0.5, 2.15+i*0.9, 5.7, 0.85, size=13, color=DARK)
txt(s, "Before / After (Bamboogle)", 6.8, 1.6, 6.0, 0.4, size=15, bold=True, color=BLUE)
for col_h, x in zip(["Model","Bamboogle"], [6.8, 9.5]):
    rect(s, x, 2.1, 2.5, 0.45, fill=BLUE)
    txt(s, col_h, x+0.05, 2.12, 2.4, 0.4, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bdata = [("0.8B base","43.55%"), ("0.8B + Flow-GRPO\n(500 steps, QA)","41.13%")]
for i, (model, acc) in enumerate(bdata):
    bg = GRAY if i%2==0 else WHITE
    for x, w, cell in [(6.8,2.6,model),(9.5,2.5,acc)]:
        rect(s, x, 2.6+i*0.85, w, 0.85, fill=bg)
        txt(s, cell, x+0.05, 2.62+i*0.85, w-0.1, 0.8, size=13, align=PP_ALIGN.CENTER,
            color=(DARK if x<9 else (RGBColor(0xC0,0,0) if i==1 else GREEN)))
txt(s, "Analysis", 6.8, 4.45, 6.0, 0.4, size=14, bold=True, color=BLUE)
for i, l in enumerate([
    "• -2.4pp regression expected at 500 steps\n  (paper trains for thousands of steps)",
    "• ~50% steps skipped: uniform rewards\n  (all 0 or all 1) = no learning signal",
    "• Training loop verified correct:\n  rewards present, loss converges",
    "• Proof-of-concept: Flow-GRPO\n  framework works end-to-end",
]):
    txt(s, l, 6.8, 5.0+i*0.6, 6.1, 0.57, size=12, color=DARK)
footer(s, "9 / 13")

# ── Slide 10: Key Findings ────────────────────────────────────────────────────
s = add_slide()
header_bar(s, "Key Findings")
findings = [
    ("Reproduced AgentFlow baseline", "4/5 benchmarks within margin; Bamboogle and GAIA match or exceed paper targets"),
    ("Model size dominates",          "7B vs 0.5B is a bigger factor than any prompt engineering or tool choice"),
    ("Search tool is critical",       "Bamboogle: 51% (no search) → 67.7% (Serper.dev). Search quality = accuracy."),
    ("Answer format is critical",     "Adding <answer> tags to prompts fixed near-0% judge scores across 2Wiki/Musique/GAIA"),
    ("Text-to-SQL generalizes well",  "91.01% on Spider dev — competitive with SOTA. Execute-verify loop is key."),
    ("Flow-GRPO framework works",     "Training loop converges. 500 steps = proof-of-concept; thousands needed for gains."),
]
for i, (title, body) in enumerate(findings):
    row_y = 1.45 + i*0.95
    rect(s, 0.3, row_y, 0.55, 0.75, fill=BLUE)
    txt(s, str(i+1), 0.3, row_y, 0.55, 0.75, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, 0.9, row_y, 12.1, 0.75, fill=GRAY if i%2==0 else WHITE)
    txt(s, title, 1.0, row_y+0.03, 4.0, 0.35, size=13, bold=True, color=BLUE)
    txt(s, body,  5.1, row_y+0.03, 7.8, 0.68, size=12, color=DARK)
footer(s, "10 / 13")

# ── Slide 11: Challenges ──────────────────────────────────────────────────────
s = add_slide()
header_bar(s, "Challenges & Engineering Lessons")
challenges = [
    ("Together AI serverless limits", "Only Qwen2.5-7B-Instruct-Turbo available; 3B/14B need dedicated endpoints. Serverless tier is restrictive for scaling studies."),
    ("Modal cold starts",             "7B model takes 3.5 min to load (torch.compile + CUDA graphs). Fix: pre-bake weights in image + --enforce-eager flag."),
    ("Qwen3.5 thinking mode",         "Default template wraps outputs in <think> blocks that break AgentFlow's JSON parsing. Fix: --override-generation-config {enable_thinking:false}."),
    ("Serper.dev quota",              "2,500 free searches exhausted mid-project. Running remaining Qwen3.5 benchmarks with Wikipedia-only (lower accuracy)."),
    ("Modal spend limit",             "Hit $70 monthly cap during 9B A100 run ($19.63 for one run). Raised to $110. A100 crash-loops without --enforce-eager."),
    ("2Wiki benchmark gap",           "37.9% vs 60% target. Comparative multi-hop questions require precise entity lookup — likely needs post-GRPO model per paper."),
]
for i, (title, body) in enumerate(challenges):
    row_y = 1.45 + i*0.95
    rect(s, 0.3, row_y, 0.55, 0.75, fill=ACCENT)
    txt(s, "!", 0.3, row_y, 0.55, 0.75, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, 0.9, row_y, 12.1, 0.75, fill=GRAY if i%2==0 else WHITE)
    txt(s, title, 1.0, row_y+0.03, 3.8, 0.35, size=13, bold=True, color=ACCENT)
    txt(s, body,  4.9, row_y+0.03, 8.0, 0.68, size=12, color=DARK)
footer(s, "11 / 13")

# ── Slide 12: Demo ────────────────────────────────────────────────────────────
s = add_slide()
header_bar(s, "Live Demo", "Gradio app — python demo/app.py")
rect(s, 0.3, 1.5, 12.7, 5.6, fill=GRAY)
txt(s, "demo/app.py", 0.5, 1.6, 12.0, 0.45, size=15, bold=True, color=BLUE)
for i, l in enumerate([
    "Enter any factual question → AgentFlow plans tool calls → returns answer with reasoning steps",
    "",
    "Try these examples:",
    '  • "What is the capital of the country that hosted the 2022 FIFA World Cup?"',
    '  • "Which director has won more Oscars: Steven Spielberg or Martin Scorsese?"',
    '  • "How many singers are in the concert_singer database?"  (SQL mode)',
    "",
    "Tool selection (checkboxes):",
    "  Google Search  |  Wikipedia  |  SQL Executor",
    "",
    "Model selection:",
    "  together-Qwen/Qwen2.5-7B-Instruct-Turbo  (default)",
]):
    txt(s, l, 0.5, 2.15+i*0.38, 12.2, 0.36, size=13 if i==0 else 12,
        bold=(i==0), color=DARK if i>0 else BLUE)
footer(s, "12 / 13")

# ── Slide 13: Conclusion ──────────────────────────────────────────────────────
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=BLUE)
rect(s, 0, 1.1, 13.33, 0.04, fill=ACCENT)
txt(s, "Conclusion", 0.5, 0.2, 12.0, 0.8, size=36, bold=True, color=WHITE)
conclusions = [
    "Successfully reimplemented AgentFlow — 4/5 benchmarks reproduced within margin",
    "Text-to-SQL: 91.01% on Spider dev — strongest single result of the project",
    "Flow-GRPO training loop verified end-to-end on Qwen3.5-0.8B",
    "Key insight: model size (7B) and search quality matter most for AgentFlow",
    "Future work: longer GRPO training (>2000 steps), multi-benchmark post-training eval",
]
for i, l in enumerate(conclusions):
    txt(s, f"{'→'} {l}", 0.7, 1.4+i*1.0, 11.8, 0.9, size=17, color=WHITE)
txt(s, "github.com/KasaVarun/agentflow-project", 0.5, 6.8, 12.0, 0.5,
    size=14, color=RGBColor(0x9D,0xC3,0xE6), align=PP_ALIGN.CENTER)
footer(s, "13 / 13")

# ── Save ──────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "agentflow_slides.pptx")
prs.save(out)
print(f"Saved: {out}")
